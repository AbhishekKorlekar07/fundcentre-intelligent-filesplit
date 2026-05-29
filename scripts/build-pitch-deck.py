"""Generate the FundCentre Intelligent Filesplit pitch deck (.pptx).

Minimalist version: each slide has a title + one short lead line.
Speaker narrates the rest from the demo script.

Usage: python3 scripts/build-pitch-deck.py
Output: pitch/FundCentre-Intelligent-Filesplit.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand palette (matches Tailwind brand-* used in the app)
BRAND_900 = RGBColor(0x0B, 0x2A, 0x4A)   # deep navy
BRAND_700 = RGBColor(0x12, 0x4A, 0x83)
BRAND_500 = RGBColor(0x2E, 0x7B, 0xC8)
INK = RGBColor(0x14, 0x1A, 0x24)
SUBTLE = RGBColor(0x55, 0x5F, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM_WHITE = RGBColor(0xCB, 0xD9, 0xEC)

OUT_DIR = "pitch"
OUT_FILE = os.path.join(OUT_DIR, "FundCentre-Intelligent-Filesplit.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


# --- helpers ---------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_footer(slide, page_num, total):
    add_rect(slide, 0, prs.slide_height - Inches(0.35),
             prs.slide_width, Inches(0.35), BRAND_900)
    add_text(slide, Inches(0.5), prs.slide_height - Inches(0.35),
             Inches(8), Inches(0.35),
             "FundCentre Intelligent Filesplit  ·  Hackathon 2026",
             size=10, color=WHITE)
    add_text(slide, prs.slide_width - Inches(1.5),
             prs.slide_height - Inches(0.35), Inches(1), Inches(0.35),
             f"{page_num} / {total}",
             size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def content_slide(eyebrow, title, lead):
    """A standard content slide: eyebrow + big title + one short lead line."""
    s = prs.slides.add_slide(BLANK)
    # left accent bar
    add_rect(s, Inches(0.6), Inches(2.4), Inches(0.15), Inches(2.0), BRAND_500)
    # eyebrow
    add_text(s, Inches(0.95), Inches(2.4), Inches(11.5), Inches(0.4),
             eyebrow.upper(), size=14, bold=True, color=BRAND_500)
    # main title
    add_text(s, Inches(0.95), Inches(2.85), Inches(11.5), Inches(1.4),
             title, size=44, bold=True, color=INK)
    # one short lead line
    add_text(s, Inches(0.95), Inches(4.5), Inches(11.5), Inches(0.8),
             lead, size=20, color=SUBTLE)
    return s


def hero_slide(eyebrow, title, lead):
    """Dark full-bleed slide for title + demo + close."""
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, BRAND_900)
    add_rect(s, 0, Inches(3.4), prs.slide_width, Inches(0.06), BRAND_500)
    add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(0.5),
             eyebrow.upper(), size=16, bold=True, color=BRAND_500)
    add_text(s, Inches(0.8), Inches(3.55), Inches(12), Inches(1.4),
             title, size=54, bold=True, color=WHITE)
    if lead:
        add_text(s, Inches(0.8), Inches(4.7), Inches(12), Inches(0.8),
                 lead, size=22, color=DIM_WHITE)
    return s


# --- slides ----------------------------------------------------------------

TOTAL = 9


def slide1_title():
    s = hero_slide("Hackathon 2026",
                   "FundCentre Intelligent Filesplit",
                   "AI-powered document splitting and tagging for GP reporting")
    add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.4),
             "Abhishek Korlekar  ·  SS&C Intralinks",
             size=14, color=DIM_WHITE)


def slide2_pain():
    s = content_slide("The problem",
                      "GPs tag every page by hand.",
                      "200+ pages. Three external IDs typed per page. One typo misroutes a capital call.")
    add_footer(s, 2, TOTAL)


def slide3_solution():
    s = content_slide("Our solution",
                      "Upload → AI reads → Review → Done.",
                      "No splitter codes. No templates. The AI reads each page like a human.")
    add_footer(s, 3, TOTAL)


def slide4_demo():
    hero_slide("Demo", "Live demo — 4 minutes",
               "Combined PDF → AI extraction → Review with preview → Split ZIP")


def slide5_how():
    s = content_slide("How it works",
                      "Three problems. One pipeline.",
                      "Multiple names per page · continuation pages · hallucinations — all handled.")
    add_footer(s, 5, TOTAL)


def slide6_innovation():
    s = content_slide("Innovation",
                      "Zero markers. Validated AI.",
                      "Most “AI document tools” stop at extraction. We close the loop.")
    add_footer(s, 6, TOTAL)


def slide7_impact():
    s = content_slide("Business impact",
                      "100 minutes → 3 minutes.",
                      "Per 200-page combined PDF. 30–50× throughput at the same headcount.")
    add_footer(s, 7, TOTAL)


def slide8_roadmap():
    s = content_slide("Roadmap",
                      "Pilot → learn → scale.",
                      "Multi-investor pages, learn-from-corrections, public API for GP pipelines.")
    add_footer(s, 8, TOTAL)


def slide9_close():
    s = hero_slide("Ask",
                   "Pilot with one GP next quarter.",
                   "Working prototype. Stays in SS&C’s AWS tenancy via Bedrock.")
    add_text(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.4),
             "Thank you.  ·  Q&A",
             size=18, bold=True, color=DIM_WHITE)


# --- build -----------------------------------------------------------------

slide1_title()
slide2_pain()
slide3_solution()
slide4_demo()
slide5_how()
slide6_innovation()
slide7_impact()
slide8_roadmap()
slide9_close()

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT_FILE)
print(f"wrote {OUT_FILE}  ({os.path.getsize(OUT_FILE) // 1024} KB)")
